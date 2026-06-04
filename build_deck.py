from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
BG       = RGBColor(0x0F, 0x0F, 0x10)   # near-black
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xF0, 0xF0, 0xF0)
SILVER   = RGBColor(0xC0, 0xC0, 0xC8)
BLUE     = RGBColor(0x3D, 0x9B, 0xFF)   # electric blue
DARK2    = RGBColor(0x16, 0x16, 0x18)
DARK3    = RGBColor(0x1E, 0x1E, 0x22)

# ── Slide dimensions ─────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank


def add_slide():
    return prs.slides.add_slide(blank_layout)


def bg(slide, color=BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def txbox(slide, text, left, top, width, height,
          font_name="Inter", font_size=18, bold=False, italic=False,
          color=WHITE, align=PP_ALIGN.LEFT, line_space=None):
    """Add a simple text box and return the shape."""
    tf = slide.shapes.add_textbox(left, top, width, height)
    frame = tf.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = align
    if line_space:
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_space * 100)))
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf


def multiline_box(slide, lines, left, top, width, height,
                  default_font="Inter", default_size=18,
                  default_color=WHITE, align=PP_ALIGN.LEFT,
                  default_bold=False):
    """
    lines = list of dicts:
      { text, font, size, color, bold, italic, align, space_before }
    """
    tf = slide.shapes.add_textbox(left, top, width, height)
    frame = tf.text_frame
    frame.word_wrap = True

    first = True
    for spec in lines:
        if first:
            p = frame.paragraphs[0]
            first = False
        else:
            p = frame.add_paragraph()

        p.alignment = spec.get("align", align)
        if spec.get("space_before"):
            p.space_before = Pt(spec["space_before"])

        run = p.add_run()
        run.text = spec.get("text", "")
        run.font.name  = spec.get("font",  default_font)
        run.font.size  = Pt(spec.get("size",  default_size))
        run.font.bold  = spec.get("bold",  default_bold)
        run.font.italic = spec.get("italic", False)
        run.font.color.rgb = spec.get("color", default_color)

    return tf


def rect(slide, left, top, width, height, fill_color=DARK2,
         line_color=None, line_width=Pt(1), radius=None):
    """Add a filled rectangle (card)."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def hline(slide, left, top, width, color=SILVER, alpha=None):
    """Draw a 1-pt horizontal rule."""
    from pptx.util import Pt
    line = slide.shapes.add_shape(1, left, top, width, Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — OPENING TITLE
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

# Background gradient bar (simulate with a dark rect at top)
top_bar = rect(s, 0, 0, W, H, fill_color=RGBColor(0x10, 0x10, 0x14))
top_bar.line.fill.background()

# Glow accent rect (center, very subtle blue tint rectangle)
glow = rect(s, Inches(3), Inches(1.5), Inches(7.33), Inches(5),
            fill_color=RGBColor(0x12, 0x18, 0x28))
glow.line.fill.background()

# Eyebrow
txbox(s, "JUNE 4TH  ·  THE RING WORKSPACE  ·  CLEARWATER, FL",
      Inches(1), Inches(1.4), Inches(11.33), Inches(0.4),
      font_name="Montserrat", font_size=11, color=SILVER,
      align=PP_ALIGN.CENTER)

# Rule
hline(s, Inches(5.5), Inches(1.95), Inches(2.33), color=SILVER)

# Main wordmark
txbox(s, "LINK'D UP",
      Inches(1), Inches(2.1), Inches(11.33), Inches(2.2),
      font_name="Arial Black", font_size=100, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)

# Tagline
multiline_box(s, [
    {"text": "Creators.  Entrepreneurs.  Collaborators.",
     "font": "Montserrat", "size": 22, "color": OFFWHITE,
     "bold": False, "align": PP_ALIGN.CENTER},
], Inches(1), Inches(4.5), Inches(11.33), Inches(1.0))

# Rule 2
hline(s, Inches(5.5), Inches(5.6), Inches(2.33), color=SILVER)

# Date/time line
txbox(s, "7:00 PM  ·  DOORS OPEN",
      Inches(1), Inches(5.75), Inches(11.33), Inches(0.4),
      font_name="Montserrat", font_size=12, color=SILVER,
      align=PP_ALIGN.CENTER)

# Bottom event bar
bar = rect(s, 0, Inches(7.1), W, Inches(0.4),
           fill_color=RGBColor(0x16, 0x16, 0x18))
txbox(s, "LINK'D UP  ·  JUNE 4TH  ·  7:00 PM  ·  THE RING WORKSPACE  ·  CLEARWATER, FL",
      Inches(0.3), Inches(7.12), Inches(12.73), Inches(0.3),
      font_name="Montserrat", font_size=9, color=SILVER,
      align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — WHY LINK'D UP EXISTS
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

# Right photo placeholder panel
photo_panel = rect(s, Inches(6.8), 0, Inches(6.53), H,
                   fill_color=RGBColor(0x1A, 0x1A, 0x20))
# Gradient fade suggestion rect
fade = rect(s, Inches(5.5), 0, Inches(2.5), H,
            fill_color=RGBColor(0x0F, 0x0F, 0x10))
fade.line.fill.background()

# Photo label
txbox(s, "[ PHOTO: Candid networking moment — warm, golden light ]",
      Inches(7.0), Inches(3.3), Inches(5.8), Inches(0.6),
      font_name="Montserrat", font_size=10, color=SILVER,
      align=PP_ALIGN.CENTER, italic=True)

# Left text column
txbox(s, "WHY WE'RE HERE",
      Inches(0.6), Inches(0.7), Inches(6.0), Inches(0.35),
      font_name="Montserrat", font_size=11, color=SILVER)

multiline_box(s, [
    {"text": "THE WORLD", "font": "Arial Black", "size": 68,
     "color": WHITE, "bold": True},
    {"text": "IS LOUD.", "font": "Arial Black", "size": 68,
     "color": WHITE, "bold": True},
    {"text": "THE ROOM", "font": "Arial Black", "size": 68,
     "color": WHITE, "bold": True},
    {"text": "IS REAL.", "font": "Arial Black", "size": 68,
     "color": WHITE, "bold": True},
], Inches(0.6), Inches(1.1), Inches(6.0), Inches(3.5))

hline(s, Inches(0.6), Inches(4.75), Inches(5.5), color=SILVER)

multiline_box(s, [
    {"text": "Social media gives you followers.",
     "font": "Montserrat", "size": 15, "color": OFFWHITE, "space_before": 6},
    {"text": "This room gives you collaborators.",
     "font": "Montserrat", "size": 15, "color": OFFWHITE, "space_before": 4},
    {"text": " ", "font": "Montserrat", "size": 8, "color": OFFWHITE},
    {"text": "The greatest opportunities in your career",
     "font": "Montserrat", "size": 15, "color": OFFWHITE, "space_before": 4},
    {"text": "won't come from an algorithm.",
     "font": "Montserrat", "size": 15, "color": OFFWHITE},
    {"text": "They'll come from a conversation",
     "font": "Montserrat", "size": 15, "color": OFFWHITE},
    {"text": "that starts in a room like this one.",
     "font": "Montserrat", "size": 15, "color": OFFWHITE},
    {"text": " ", "font": "Montserrat", "size": 8, "color": OFFWHITE},
    {"text": "Great things happen in person.",
     "font": "Arial Black", "size": 22, "color": BLUE, "bold": True},
], Inches(0.6), Inches(4.9), Inches(6.0), Inches(2.4))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT IS LINK'D UP?
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

# Subtle texture rect
rect(s, 0, 0, W, H, fill_color=RGBColor(0x12, 0x12, 0x16))

txbox(s, "THE MISSION",
      Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.35),
      font_name="Montserrat", font_size=11, color=SILVER,
      align=PP_ALIGN.CENTER)

multiline_box(s, [
    {"text": "NOT YOUR AVERAGE", "font": "Arial Black", "size": 60,
     "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER},
    {"text": "NETWORKING EVENT.", "font": "Arial Black", "size": 60,
     "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER},
], Inches(0.5), Inches(0.8), Inches(12.33), Inches(1.6))

hline(s, Inches(5.7), Inches(2.55), Inches(2.0), color=SILVER)

multiline_box(s, [
    {"text": "Link'd Up is Tampa Bay's creative community",
     "font": "Montserrat", "size": 15, "color": OFFWHITE,
     "align": PP_ALIGN.CENTER},
    {"text": "for the people who make things happen —",
     "font": "Montserrat", "size": 15, "color": OFFWHITE,
     "align": PP_ALIGN.CENTER},
    {"text": "the photographers, founders, influencers, designers,",
     "font": "Montserrat", "size": 15, "color": OFFWHITE,
     "align": PP_ALIGN.CENTER},
    {"text": "directors, and dreamers ready to build something bigger.",
     "font": "Montserrat", "size": 15, "color": OFFWHITE,
     "align": PP_ALIGN.CENTER},
    {"text": " ", "size": 6, "color": OFFWHITE},
    {"text": "This is a social event. A creative exchange.",
     "font": "Montserrat", "size": 15, "color": OFFWHITE,
     "align": PP_ALIGN.CENTER, "bold": True},
    {"text": "A place where business feels like a conversation, not a pitch.",
     "font": "Montserrat", "size": 15, "color": OFFWHITE,
     "align": PP_ALIGN.CENTER},
], Inches(0.8), Inches(2.7), Inches(11.73), Inches(1.8))

# Three pillar cards
card_w = Inches(4.0)
card_h = Inches(2.6)
card_y = Inches(4.65)
labels = ["WHO IT'S FOR", "WHAT WE DO", "HOW WE'RE DIFFERENT"]
titles = ["Creators & Builders", "Connect & Collaborate", "Social, Not Corporate"]
bodies = [
    "Photographers. Videographers.\nDesigners. Models. Entrepreneurs.\nAgencies. Artists. Influencers.\nAnyone building something real.",
    "We curate genuine connections.\nNot business cards. Not pitches.\nReal conversations. Real partnerships.\nReal opportunities.",
    "No awkward icebreakers.\nNo name-tag circles.\nJust good people in a great space,\nmaking moves together.",
]
gaps = [Inches(0.22), Inches(4.55), Inches(8.88)]

for i in range(3):
    rect(s, gaps[i], card_y, card_w, card_h,
         fill_color=DARK2, line_color=SILVER)
    txbox(s, labels[i],
          gaps[i]+Inches(0.2), card_y+Inches(0.18), card_w-Inches(0.4), Inches(0.3),
          font_name="Montserrat", font_size=9, color=BLUE)
    txbox(s, titles[i],
          gaps[i]+Inches(0.2), card_y+Inches(0.52), card_w-Inches(0.4), Inches(0.5),
          font_name="Arial Black", font_size=18, bold=True, color=WHITE)
    txbox(s, bodies[i],
          gaps[i]+Inches(0.2), card_y+Inches(1.05), card_w-Inches(0.4), Inches(1.4),
          font_name="Montserrat", font_size=11.5, color=SILVER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE VISION
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

txbox(s, "WHERE WE'RE GOING",
      Inches(0.5), Inches(0.35), Inches(12.33), Inches(0.32),
      font_name="Montserrat", font_size=11, color=SILVER,
      align=PP_ALIGN.CENTER)

txbox(s, "BIGGER. EVERY. TIME.",
      Inches(0.5), Inches(0.65), Inches(12.33), Inches(1.2),
      font_name="Arial Black", font_size=56, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)

hline(s, Inches(5.7), Inches(1.85), Inches(2.0), color=SILVER)

txbox(s, "Link'd Up is built to grow with Tampa Bay's creative economy.",
      Inches(1.0), Inches(2.0), Inches(11.33), Inches(0.4),
      font_name="Montserrat", font_size=13, color=OFFWHITE,
      align=PP_ALIGN.CENTER)

nums = ["01","02","03","04","05","06"]
vtitles = ["Monthly Events","Tampa Bay Expansion","Creator Spotlights",
           "Brand Collaborations","Educational Workshops","Community Platform"]
vbodies = [
    "Link'd Up runs every month —\nconsistent energy, new faces,\nnew opportunities every time.",
    "Starting in Clearwater, growing\ninto Tampa, St. Pete, and beyond.\nThe whole Bay, connected.",
    "Every event spotlights a creator\nfrom the community — their story,\ntheir work, their journey.",
    "Building partnerships with local\nbrands, agencies, and businesses\nwho invest in creatives.",
    "Future events include workshops\non content creation, brand building,\nand business development.",
    "A private community where members\nstay connected, share work, find\ncollaborators, and grow.",
]

cw = Inches(3.9)
ch = Inches(1.95)
cy1 = Inches(2.6)
cy2 = Inches(4.65)
cx = [Inches(0.22), Inches(4.55), Inches(8.88)]

for i in range(6):
    row = i // 3
    col = i % 3
    x = cx[col]
    y = cy1 if row == 0 else cy2
    rect(s, x, y, cw, ch, fill_color=DARK2, line_color=RGBColor(0x3D,0x9B,0xFF))
    txbox(s, nums[i],
          x+Inches(0.15), y+Inches(0.12), Inches(0.6), Inches(0.55),
          font_name="Arial Black", font_size=26, bold=True, color=BLUE)
    txbox(s, vtitles[i],
          x+Inches(0.15), y+Inches(0.65), cw-Inches(0.3), Inches(0.38),
          font_name="Arial Black", font_size=14, bold=True, color=WHITE)
    txbox(s, vbodies[i],
          x+Inches(0.15), y+Inches(1.05), cw-Inches(0.3), Inches(0.82),
          font_name="Montserrat", font_size=10, color=SILVER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TONIGHT'S EXPERIENCE (Timeline)
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

txbox(s, "TONIGHT'S EXPERIENCE",
      Inches(0.5), Inches(0.35), Inches(12.33), Inches(0.32),
      font_name="Montserrat", font_size=11, color=SILVER,
      align=PP_ALIGN.CENTER)

txbox(s, "HERE'S HOW THE NIGHT UNFOLDS.",
      Inches(0.5), Inches(0.65), Inches(12.33), Inches(1.1),
      font_name="Arial Black", font_size=48, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)

hline(s, Inches(5.7), Inches(1.75), Inches(2.0), color=SILVER)

# Timeline track
track_y = Inches(3.5)
hline(s, Inches(0.8), track_y, Inches(11.73), color=SILVER)

times  = ["7:00 PM","7:15 PM","7:40 PM","8:20 PM","8:35 PM"]
mtitles = ["Doors Open\n& Mingle","Welcome &\nIntroductions",
           "LINK UP\nHOUR","Content\nCorner","Closing &\nNext Event"]
details = [
    "Arrive, check in,\nfind your vibe.",
    "Host intro.\nWho's in the room.",
    "Meet people.\nExchange contacts.",
    "Photos. Videos.\nCreator spotlights.",
    "Final words.\nNext event revealed.",
]

node_xs = [Inches(1.1), Inches(3.55), Inches(6.0), Inches(8.45), Inches(10.9)]

for i, nx in enumerate(node_xs):
    is_featured = (i == 2)
    nc = BLUE
    # Node circle (simulate with small rect)
    dot = rect(s, nx - Inches(0.12), track_y - Inches(0.12),
               Inches(0.24), Inches(0.24),
               fill_color=(BLUE if is_featured else RGBColor(0x0F,0x0F,0x10)),
               line_color=BLUE)

    # Time above
    txbox(s, times[i],
          nx - Inches(0.7), track_y - Inches(0.65), Inches(1.4), Inches(0.45),
          font_name="Arial Black", font_size=16, bold=True,
          color=(BLUE if is_featured else SILVER), align=PP_ALIGN.CENTER)

    # Card below
    card_top = track_y + Inches(0.25)
    cr = rect(s, nx - Inches(0.95), card_top, Inches(1.9), Inches(2.5),
              fill_color=(RGBColor(0x10,0x18,0x2A) if is_featured else DARK2),
              line_color=(BLUE if is_featured else RGBColor(0x50,0x50,0x60)))
    txbox(s, mtitles[i],
          nx - Inches(0.85), card_top + Inches(0.15),
          Inches(1.7), Inches(0.75),
          font_name="Arial Black", font_size=13, bold=True,
          color=(BLUE if is_featured else WHITE), align=PP_ALIGN.CENTER)
    txbox(s, details[i],
          nx - Inches(0.85), card_top + Inches(1.0),
          Inches(1.7), Inches(1.3),
          font_name="Montserrat", font_size=10, color=SILVER,
          align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MEET THE COMMUNITY
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

txbox(s, "THIS IS YOUR COMMUNITY",
      Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.32),
      font_name="Montserrat", font_size=11, color=SILVER,
      align=PP_ALIGN.CENTER)

txbox(s, "LOOK AROUND. THIS IS YOUR ROOM.",
      Inches(0.5), Inches(0.6), Inches(12.33), Inches(1.0),
      font_name="Arial Black", font_size=44, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)

community = [
    ("CREATORS","Making culture move."),
    ("ENTREPRENEURS","Building what's next."),
    ("PHOTOGRAPHERS","Capturing real moments."),
    ("MODELS","Setting the visual standard."),
    ("ARTISTS","The original disruptors."),
    ("VIDEOGRAPHERS","Story through the lens."),
    ("DESIGNERS","Shaping how the world looks."),
    ("AGENCY OWNERS","Brands behind the brands."),
    ("CONTENT CREATORS","The new media."),
]

cols = 3
card_w = Inches(4.0)
card_h = Inches(1.75)
gap_x  = Inches(0.22)
gap_y  = Inches(0.18)
start_x = Inches(0.22)
start_y = Inches(1.75)

for idx, (label, sub) in enumerate(community):
    col = idx % cols
    row = idx // cols
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)
    # Card background — alternate subtle shade
    shade = RGBColor(0x14,0x14,0x18) if (row + col) % 2 == 0 else DARK2
    rect(s, x, y, card_w, card_h, fill_color=shade,
         line_color=RGBColor(0x40,0x40,0x50))
    # Label
    txbox(s, label,
          x+Inches(0.15), y+Inches(0.35), card_w-Inches(0.3), Inches(0.6),
          font_name="Arial Black", font_size=18, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, sub,
          x+Inches(0.15), y+Inches(1.1), card_w-Inches(0.3), Inches(0.45),
          font_name="Montserrat", font_size=10, color=SILVER,
          align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — HOW TO GET THE MOST OUT OF TONIGHT
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

txbox(s, "YOUR PLAYBOOK FOR TONIGHT",
      Inches(0.6), Inches(0.35), Inches(12.0), Inches(0.32),
      font_name="Montserrat", font_size=11, color=SILVER,
      align=PP_ALIGN.CENTER)

txbox(s, "HOW TO WIN THIS ROOM.",
      Inches(0.6), Inches(0.65), Inches(12.0), Inches(1.1),
      font_name="Arial Black", font_size=52, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)

hline(s, Inches(5.7), Inches(1.75), Inches(2.0), color=SILVER)

step_nums  = ["01","02","03","04","05"]
step_title = [
    "Introduce Yourself First.",
    "Ask Better Questions.",
    "Exchange Your Socials.",
    "Look for the Overlap.",
    "Follow Up Tomorrow.",
]
step_body = [
    "Walk up. Say your name. Say what you do in one sentence. That's all it takes.",
    "Don't ask 'what do you do.' Ask what they're working on. That's how you find alignment.",
    "Instagram is the new business card. Follow each other before you leave the conversation.",
    "What do they need that you have? What do you need that they have? That's where collabs live.",
    "The meeting is tonight. The relationship is built tomorrow. Send the DM. Lock in the coffee.",
]

row_h  = Inches(0.95)
start_y = Inches(2.0)
for i in range(5):
    y = start_y + i * row_h
    # Big background number
    txbox(s, step_nums[i],
          Inches(0.3), y - Inches(0.2), Inches(1.2), Inches(1.0),
          font_name="Arial Black", font_size=52, bold=True,
          color=RGBColor(0x3D,0x9B,0xFF), align=PP_ALIGN.LEFT)
    # Title
    txbox(s, step_title[i],
          Inches(1.7), y, Inches(5.0), Inches(0.4),
          font_name="Arial Black", font_size=20, bold=True, color=WHITE)
    # Body
    txbox(s, step_body[i],
          Inches(1.7), y + Inches(0.42), Inches(10.7), Inches(0.45),
          font_name="Montserrat", font_size=12.5, color=SILVER)
    if i < 4:
        hline(s, Inches(1.7), y + Inches(0.9), Inches(10.7),
              color=RGBColor(0x30,0x30,0x38))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — LINK UP HOUR
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

# Deep blue glow bg
rect(s, Inches(2), Inches(0.5), Inches(9.33), Inches(6.5),
     fill_color=RGBColor(0x0A, 0x12, 0x22))

txbox(s, "7:40 PM  —  THIS IS IT",
      Inches(0.5), Inches(0.35), Inches(12.33), Inches(0.35),
      font_name="Montserrat", font_size=12, color=SILVER,
      align=PP_ALIGN.CENTER)

multiline_box(s, [
    {"text": "LINK", "font": "Arial Black", "size": 130,
     "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER},
    {"text": "UP", "font": "Arial Black", "size": 130,
     "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER},
    {"text": "HOUR.", "font": "Arial Black", "size": 130,
     "color": BLUE, "bold": True, "align": PP_ALIGN.CENTER},
], Inches(0.5), Inches(0.6), Inches(12.33), Inches(4.3))

hline(s, Inches(4.5), Inches(4.85), Inches(4.33), color=SILVER)

# Four callout words
words  = ["MEET","DISCOVER","BUILD","CREATE"]
subs   = ["NEW PEOPLE","OPPORTUNITIES","RELATIONSHIPS","PARTNERSHIPS"]
word_x = [Inches(0.5), Inches(3.7), Inches(6.9), Inches(10.1)]

for i in range(4):
    txbox(s, words[i],
          word_x[i], Inches(5.05), Inches(3.0), Inches(0.6),
          font_name="Arial Black", font_size=28, bold=True,
          color=BLUE, align=PP_ALIGN.CENTER)
    txbox(s, subs[i],
          word_x[i], Inches(5.65), Inches(3.0), Inches(0.3),
          font_name="Montserrat", font_size=9, color=SILVER,
          align=PP_ALIGN.CENTER)

txbox(s, "Move freely. Introduce yourself. Exchange socials. Find who you've been looking for.",
      Inches(0.5), Inches(6.25), Inches(12.33), Inches(0.45),
      font_name="Montserrat", font_size=13, color=OFFWHITE,
      align=PP_ALIGN.CENTER)

txbox(s, "GO LINK UP.",
      Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.5),
      font_name="Arial Black", font_size=22, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CONTENT CORNER
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

# Right photo panel
rect(s, Inches(6.8), 0, Inches(6.53), H,
     fill_color=RGBColor(0x18, 0x18, 0x22))
txbox(s, "[ PHOTO: Creator/photographer at work — ring light, camera, studio vibe ]",
      Inches(7.0), Inches(3.1), Inches(5.8), Inches(0.7),
      font_name="Montserrat", font_size=10, color=SILVER,
      align=PP_ALIGN.CENTER, italic=True)
fade2 = rect(s, Inches(5.5), 0, Inches(2.0), H,
             fill_color=RGBColor(0x0F,0x0F,0x10))
fade2.line.fill.background()

txbox(s, "8:20 PM",
      Inches(0.6), Inches(0.55), Inches(5.5), Inches(0.32),
      font_name="Montserrat", font_size=11, color=BLUE)

multiline_box(s, [
    {"text": "CONTENT", "font": "Arial Black", "size": 72,
     "color": WHITE, "bold": True},
    {"text": "CORNER.", "font": "Arial Black", "size": 72,
     "color": WHITE, "bold": True},
], Inches(0.6), Inches(0.85), Inches(5.8), Inches(2.1))

hline(s, Inches(0.6), Inches(3.1), Inches(5.5), color=SILVER)

txbox(s, "This is your moment to create.\nBring the energy. Be yourself.\nThe content you make tonight\nbuilds your brand tomorrow.",
      Inches(0.6), Inches(3.25), Inches(5.5), Inches(1.3),
      font_name="Montserrat", font_size=14, color=OFFWHITE)

# Content type chips grid
chips = ["PHOTOS","VIDEOS","CREATOR SPOTLIGHTS","INTERVIEWS","GROUP SHOTS","BEHIND THE SCENES"]
chip_w = Inches(2.1)
chip_h = Inches(0.38)
chip_x = [Inches(0.6), Inches(2.9), Inches(5.2)]
chip_y = [Inches(4.65), Inches(5.12)]

for i, chip in enumerate(chips):
    col = i % 3
    row = i // 3
    cx2 = chip_x[col]
    cy2 = chip_y[row]
    rect(s, cx2, cy2, chip_w, chip_h,
         fill_color=DARK3, line_color=SILVER)
    txbox(s, chip, cx2 + Inches(0.1), cy2 + Inches(0.07),
          chip_w - Inches(0.2), chip_h - Inches(0.1),
          font_name="Montserrat", font_size=9.5, color=WHITE,
          align=PP_ALIGN.CENTER)

hline(s, Inches(0.6), Inches(5.62), Inches(5.5), color=SILVER)

txbox(s, "STEP IN. GET SHOT. GET FEATURED.",
      Inches(0.6), Inches(5.75), Inches(5.5), Inches(0.45),
      font_name="Arial Black", font_size=17, bold=True, color=BLUE)

txbox(s, "Tag @linkdup.co · #linkdup · #linkdupclearwater",
      Inches(0.6), Inches(6.22), Inches(5.5), Inches(0.35),
      font_name="Montserrat", font_size=12, color=SILVER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — THE FUTURE OF LINK'D UP
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

txbox(s, "THE ROAD AHEAD",
      Inches(0.6), Inches(0.4), Inches(6.0), Inches(0.32),
      font_name="Montserrat", font_size=11, color=SILVER)

multiline_box(s, [
    {"text": "TAMPA BAY'S", "font": "Arial Black", "size": 56,
     "color": WHITE, "bold": True},
    {"text": "CREATIVE", "font": "Arial Black", "size": 56,
     "color": WHITE, "bold": True},
    {"text": "COMMUNITY", "font": "Arial Black", "size": 56,
     "color": WHITE, "bold": True},
    {"text": "IS GROWING.", "font": "Arial Black", "size": 56,
     "color": BLUE, "bold": True},
], Inches(0.6), Inches(0.75), Inches(6.5), Inches(3.5))

hline(s, Inches(0.6), Inches(4.4), Inches(5.5), color=SILVER)

txbox(s, "Link'd Up starts here — tonight — in Clearwater.\nBut this is bigger than one city.\nWe're building the network that connects\nevery creative in Tampa Bay.",
      Inches(0.6), Inches(4.55), Inches(5.5), Inches(1.5),
      font_name="Montserrat", font_size=13.5, color=OFFWHITE)

# Map placeholder
map_panel = rect(s, Inches(7.2), Inches(0.6), Inches(5.7), Inches(6.1),
                 fill_color=RGBColor(0x12,0x14,0x20),
                 line_color=RGBColor(0x30,0x30,0x45))
txbox(s, "[ MAP: Tampa Bay Area\nClearwater · Tampa · St. Pete ]",
      Inches(7.4), Inches(2.8), Inches(5.3), Inches(0.8),
      font_name="Montserrat", font_size=11, color=SILVER,
      align=PP_ALIGN.CENTER, italic=True)

# City markers
cities  = ["CLEARWATER","TAMPA","ST. PETE"]
status  = ["ACTIVE · JUNE 2026","COMING SOON","COMING SOON"]
city_x  = [Inches(7.5), Inches(9.2), Inches(10.9)]
city_y  = Inches(4.8)

for i, city in enumerate(cities):
    is_active = (i == 0)
    rect(s, city_x[i], city_y, Inches(1.6), Inches(1.35),
         fill_color=(RGBColor(0x0D,0x18,0x30) if is_active else DARK2),
         line_color=(BLUE if is_active else RGBColor(0x35,0x35,0x45)))
    txbox(s, city,
          city_x[i]+Inches(0.08), city_y+Inches(0.18), Inches(1.45), Inches(0.45),
          font_name="Arial Black", font_size=13, bold=True,
          color=(BLUE if is_active else SILVER), align=PP_ALIGN.CENTER)
    txbox(s, status[i],
          city_x[i]+Inches(0.08), city_y+Inches(0.72), Inches(1.45), Inches(0.38),
          font_name="Montserrat", font_size=8, color=SILVER,
          align=PP_ALIGN.CENTER)

txbox(s, "\"One event. One community. Three cities. The whole Bay, connected.\"",
      Inches(0.6), Inches(6.2), Inches(12.0), Inches(0.45),
      font_name="Montserrat", font_size=13, color=SILVER,
      align=PP_ALIGN.CENTER, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — JOIN THE COMMUNITY (QR)
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

# Glow behind QR
rect(s, Inches(4.5), Inches(1.5), Inches(4.33), Inches(4.5),
     fill_color=RGBColor(0x0A,0x14,0x28))

txbox(s, "STAY CONNECTED",
      Inches(0.5), Inches(0.35), Inches(12.33), Inches(0.32),
      font_name="Montserrat", font_size=11, color=SILVER,
      align=PP_ALIGN.CENTER)

txbox(s, "JOIN THE COMMUNITY.",
      Inches(0.5), Inches(0.65), Inches(12.33), Inches(1.05),
      font_name="Arial Black", font_size=56, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)

hline(s, Inches(5.7), Inches(1.72), Inches(2.0), color=SILVER)

# QR code placeholder box
qr_box = rect(s, Inches(5.42), Inches(1.85), Inches(2.5), Inches(2.5),
              fill_color=WHITE, line_color=BLUE)
txbox(s, "QR CODE\nPLACEHOLDER\n\nReplace with\nlive QR link",
      Inches(5.52), Inches(2.15), Inches(2.3), Inches(1.9),
      font_name="Montserrat", font_size=11, color=RGBColor(0x33,0x33,0x33),
      align=PP_ALIGN.CENTER)

txbox(s, "Scan to stay connected",
      Inches(4.5), Inches(4.45), Inches(4.33), Inches(0.35),
      font_name="Montserrat", font_size=12, color=SILVER,
      align=PP_ALIGN.CENTER)

hline(s, Inches(2.5), Inches(4.9), Inches(8.33), color=RGBColor(0x30,0x30,0x38))

# Four link items
link_labels = ["WEBSITE","INSTAGRAM","QUESTIONNAIRE","NEXT EVENT"]
link_values = ["linkdup.co","@linkdup.co","[Link Here]","[RSVP Link]"]
lx = [Inches(0.6), Inches(3.6), Inches(6.6), Inches(9.6)]

for i in range(4):
    txbox(s, link_labels[i],
          lx[i], Inches(5.1), Inches(3.0), Inches(0.32),
          font_name="Montserrat", font_size=9, color=BLUE,
          align=PP_ALIGN.CENTER)
    txbox(s, link_values[i],
          lx[i], Inches(5.4), Inches(3.0), Inches(0.4),
          font_name="Arial Black", font_size=14, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)

hline(s, Inches(2.5), Inches(5.9), Inches(8.33), color=RGBColor(0x30,0x30,0x38))

txbox(s, "Fill out the community questionnaire — let us know who you are and what you're looking for.",
      Inches(1.0), Inches(6.05), Inches(11.33), Inches(0.4),
      font_name="Montserrat", font_size=12, color=SILVER,
      align=PP_ALIGN.CENTER)

txbox(s, "#LINKDUP  ·  #TAMPABAYCREATIVES  ·  #LINKDUPCLEARWATER",
      Inches(1.0), Inches(6.6), Inches(11.33), Inches(0.4),
      font_name="Montserrat", font_size=11, color=OFFWHITE,
      align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — THANK YOU / CLOSING
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

# Background richness
rect(s, 0, 0, W, H, fill_color=RGBColor(0x0C,0x0C,0x12))
rect(s, Inches(2), Inches(0.8), Inches(9.33), Inches(6.0),
     fill_color=RGBColor(0x0A,0x10,0x1E))

txbox(s, "THANK YOU FOR BEING HERE TONIGHT.",
      Inches(0.5), Inches(0.6), Inches(12.33), Inches(0.38),
      font_name="Montserrat", font_size=12, color=SILVER,
      align=PP_ALIGN.CENTER)

hline(s, Inches(5.0), Inches(1.08), Inches(3.33), color=SILVER)

multiline_box(s, [
    {"text": "LET'S BUILD", "font": "Arial Black", "size": 78,
     "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER},
    {"text": "SOMETHING", "font": "Arial Black", "size": 78,
     "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER},
    {"text": "REAL", "font": "Arial Black", "size": 78,
     "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER},
    {"text": "TOGETHER.", "font": "Arial Black", "size": 78,
     "color": BLUE, "bold": True, "align": PP_ALIGN.CENTER},
], Inches(0.5), Inches(1.2), Inches(12.33), Inches(4.0))

hline(s, Inches(5.0), Inches(5.3), Inches(3.33), color=SILVER)

txbox(s, "LINK'D UP",
      Inches(0.5), Inches(5.45), Inches(12.33), Inches(0.9),
      font_name="Arial Black", font_size=44, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)

txbox(s, "THE RING WORKSPACE  ·  CLEARWATER, FL  ·  JUNE 4TH",
      Inches(0.5), Inches(6.38), Inches(12.33), Inches(0.35),
      font_name="Montserrat", font_size=11, color=SILVER,
      align=PP_ALIGN.CENTER)

txbox(s, "@LINKDUP.CO  ·  #LINKDUP",
      Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.35),
      font_name="Montserrat", font_size=11, color=SILVER,
      align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# BONUS A — SPONSOR: THE RING WORKSPACE
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

rect(s, Inches(2.5), Inches(1.0), Inches(8.33), Inches(5.5),
     fill_color=RGBColor(0x12,0x12,0x18))

txbox(s, "TONIGHT'S HOME",
      Inches(0.5), Inches(0.45), Inches(12.33), Inches(0.32),
      font_name="Montserrat", font_size=11, color=SILVER,
      align=PP_ALIGN.CENTER)

txbox(s, "A SPECIAL THANK YOU TO",
      Inches(0.5), Inches(0.8), Inches(12.33), Inches(0.7),
      font_name="Arial Black", font_size=36, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)

# Logo placeholder
logo_box = rect(s, Inches(4.67), Inches(1.65), Inches(4.0), Inches(1.4),
                fill_color=RGBColor(0x1E,0x1E,0x26),
                line_color=RGBColor(0x50,0x50,0x60))
txbox(s, "THE RING WORKSPACE\n[ LOGO PLACEHOLDER ]",
      Inches(4.77), Inches(1.85), Inches(3.8), Inches(1.0),
      font_name="Arial Black", font_size=16, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)

hline(s, Inches(5.0), Inches(3.2), Inches(3.33), color=SILVER)

txbox(s, "The Ring Workspace is more than a co-working space —\nit's the creative infrastructure Tampa Bay's builders need.\nThank you for being the home of Link'd Up.",
      Inches(1.5), Inches(3.35), Inches(10.33), Inches(1.4),
      font_name="Montserrat", font_size=15, color=OFFWHITE,
      align=PP_ALIGN.CENTER)

txbox(s, "600 CLEVELAND STREET  ·  CLEARWATER, FLORIDA",
      Inches(1.5), Inches(4.85), Inches(10.33), Inches(0.35),
      font_name="Montserrat", font_size=11, color=SILVER,
      align=PP_ALIGN.CENTER)

txbox(s, "theringworkspace.com",
      Inches(1.5), Inches(5.3), Inches(10.33), Inches(0.35),
      font_name="Montserrat", font_size=12, color=BLUE,
      align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# BONUS B — NEXT EVENT TEASER
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

rect(s, Inches(1.5), Inches(0.5), Inches(10.33), Inches(6.5),
     fill_color=RGBColor(0x0A,0x10,0x1E))

txbox(s, "MARK YOUR CALENDAR",
      Inches(0.5), Inches(0.6), Inches(12.33), Inches(0.32),
      font_name="Montserrat", font_size=11, color=SILVER,
      align=PP_ALIGN.CENTER)

multiline_box(s, [
    {"text": "LINK'D UP", "font": "Arial Black", "size": 80,
     "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER},
    {"text": "IS BACK.", "font": "Arial Black", "size": 80,
     "color": BLUE, "bold": True, "align": PP_ALIGN.CENTER},
], Inches(0.5), Inches(0.95), Inches(12.33), Inches(2.4))

hline(s, Inches(4.5), Inches(3.4), Inches(4.33), color=SILVER)

txbox(s, "[ NEXT DATE · TIME — Fill in before event ]",
      Inches(0.5), Inches(3.55), Inches(12.33), Inches(0.6),
      font_name="Arial Black", font_size=26, bold=True,
      color=BLUE, align=PP_ALIGN.CENTER)

txbox(s, "THE RING WORKSPACE  ·  CLEARWATER, FL",
      Inches(0.5), Inches(4.2), Inches(12.33), Inches(0.38),
      font_name="Montserrat", font_size=12, color=SILVER,
      align=PP_ALIGN.CENTER)

# QR placeholder
qr2 = rect(s, Inches(5.67), Inches(4.7), Inches(2.0), Inches(2.0),
            fill_color=WHITE, line_color=BLUE)
txbox(s, "RSVP\nQR CODE",
      Inches(5.77), Inches(5.1), Inches(1.8), Inches(1.2),
      font_name="Montserrat", font_size=12, color=RGBColor(0x22,0x22,0x22),
      align=PP_ALIGN.CENTER)

txbox(s, "RSVP NOW  ·  SPOTS ARE LIMITED",
      Inches(0.5), Inches(6.78), Inches(12.33), Inches(0.38),
      font_name="Montserrat", font_size=11, color=SILVER,
      align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# BONUS C — EVENT FEEDBACK SLIDE
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

txbox(s, "YOUR OPINION SHAPES WHAT'S NEXT",
      Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.32),
      font_name="Montserrat", font_size=11, color=SILVER,
      align=PP_ALIGN.CENTER)

txbox(s, "HOW'D WE DO?",
      Inches(0.5), Inches(0.72), Inches(12.33), Inches(1.1),
      font_name="Arial Black", font_size=62, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)

hline(s, Inches(5.0), Inches(1.85), Inches(3.33), color=SILVER)

# Feedback QR
feedback_qr = rect(s, Inches(5.42), Inches(2.0), Inches(2.5), Inches(2.5),
                    fill_color=WHITE, line_color=BLUE)
txbox(s, "FEEDBACK\nQR CODE\n\nReplace with\nlive link",
      Inches(5.52), Inches(2.25), Inches(2.3), Inches(2.0),
      font_name="Montserrat", font_size=11, color=RGBColor(0x33,0x33,0x33),
      align=PP_ALIGN.CENTER)

txbox(s, "Scan to share your feedback — it takes 2 minutes.",
      Inches(1.0), Inches(4.65), Inches(11.33), Inches(0.38),
      font_name="Montserrat", font_size=13, color=SILVER,
      align=PP_ALIGN.CENTER)

hline(s, Inches(3.0), Inches(5.15), Inches(7.33), color=RGBColor(0x30,0x30,0x38))

txbox(s, "What you loved.  What you'd change.  Who you met.  What you're looking forward to.",
      Inches(1.0), Inches(5.3), Inches(11.33), Inches(0.4),
      font_name="Montserrat", font_size=12.5, color=OFFWHITE,
      align=PP_ALIGN.CENTER)

txbox(s, "Your feedback is how we keep this community growing in the right direction.",
      Inches(1.0), Inches(5.85), Inches(11.33), Inches(0.38),
      font_name="Montserrat", font_size=12.5, color=SILVER,
      align=PP_ALIGN.CENTER)

txbox(s, "LINK'D UP  ·  @LINKDUP.CO  ·  #LINKDUP",
      Inches(0.5), Inches(6.75), Inches(12.33), Inches(0.38),
      font_name="Montserrat", font_size=11, color=SILVER,
      align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# BONUS D — CREATOR SPOTLIGHT TEMPLATE
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

# Right photo panel
rect(s, Inches(6.5), 0, Inches(6.83), H,
     fill_color=RGBColor(0x18,0x18,0x22))
txbox(s, "[ CREATOR PHOTO\nHigh-quality portrait ]",
      Inches(7.5), Inches(3.2), Inches(4.5), Inches(0.8),
      font_name="Montserrat", font_size=11, color=SILVER,
      align=PP_ALIGN.CENTER, italic=True)
fade3 = rect(s, Inches(5.2), 0, Inches(2.0), H,
             fill_color=RGBColor(0x0F,0x0F,0x10))
fade3.line.fill.background()

txbox(s, "CREATOR SPOTLIGHT",
      Inches(0.6), Inches(0.6), Inches(5.5), Inches(0.35),
      font_name="Montserrat", font_size=11, color=BLUE)

txbox(s, "[FIRST NAME]",
      Inches(0.6), Inches(1.0), Inches(5.5), Inches(1.1),
      font_name="Arial Black", font_size=68, bold=True, color=WHITE)

txbox(s, "[LAST NAME]",
      Inches(0.6), Inches(2.0), Inches(5.5), Inches(1.1),
      font_name="Arial Black", font_size=68, bold=True, color=WHITE)

txbox(s, "[ Creator Category  ·  e.g. Photographer / Videographer ]",
      Inches(0.6), Inches(3.1), Inches(5.5), Inches(0.42),
      font_name="Arial Black", font_size=18, color=SILVER)

hline(s, Inches(0.6), Inches(3.65), Inches(5.2), color=SILVER)

txbox(s, "[ Write a 2–3 sentence bio here. Their story, their craft,\ntheir mission. Who they are. What they create. Why it matters. ]",
      Inches(0.6), Inches(3.82), Inches(5.2), Inches(1.2),
      font_name="Montserrat", font_size=13.5, color=OFFWHITE)

txbox(s, "@[HANDLE]",
      Inches(0.6), Inches(5.2), Inches(5.2), Inches(0.55),
      font_name="Arial Black", font_size=28, bold=True, color=BLUE)

txbox(s, "[ website.com ]",
      Inches(0.6), Inches(5.75), Inches(5.2), Inches(0.35),
      font_name="Montserrat", font_size=13, color=SILVER)

txbox(s, "\"[ Optional one-line quote from the creator ]\"",
      Inches(0.6), Inches(6.3), Inches(5.2), Inches(0.45),
      font_name="Montserrat", font_size=12, color=SILVER, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# BONUS E — BUSINESS SPOTLIGHT TEMPLATE
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

rect(s, Inches(1.5), Inches(0.5), Inches(10.33), Inches(6.5),
     fill_color=RGBColor(0x12,0x12,0x18))

txbox(s, "BUSINESS SPOTLIGHT",
      Inches(0.5), Inches(0.6), Inches(12.33), Inches(0.35),
      font_name="Montserrat", font_size=11, color=BLUE,
      align=PP_ALIGN.CENTER)

# Logo placeholder
rect(s, Inches(5.17), Inches(1.05), Inches(3.0), Inches(1.0),
     fill_color=RGBColor(0x1E,0x1E,0x28), line_color=RGBColor(0x45,0x45,0x55))
txbox(s, "[ BUSINESS LOGO ]",
      Inches(5.27), Inches(1.25), Inches(2.8), Inches(0.55),
      font_name="Montserrat", font_size=11, color=SILVER,
      align=PP_ALIGN.CENTER)

txbox(s, "[ BUSINESS NAME ]",
      Inches(0.5), Inches(2.2), Inches(12.33), Inches(0.9),
      font_name="Arial Black", font_size=52, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)

txbox(s, "[ Industry  ·  e.g. Creative Agency  ·  Production House ]",
      Inches(0.5), Inches(3.1), Inches(12.33), Inches(0.42),
      font_name="Arial Black", font_size=18, color=SILVER,
      align=PP_ALIGN.CENTER)

hline(s, Inches(4.0), Inches(3.65), Inches(5.33), color=SILVER)

txbox(s, "[ 2–3 sentences about what the business does,\nwho they serve, and what makes them part of this community. ]",
      Inches(1.5), Inches(3.82), Inches(10.33), Inches(0.9),
      font_name="Montserrat", font_size=13.5, color=OFFWHITE,
      align=PP_ALIGN.CENTER)

hline(s, Inches(2.5), Inches(4.85), Inches(8.33), color=RGBColor(0x30,0x30,0x38))

# Two info columns
txbox(s, "LOOKING FOR",
      Inches(2.0), Inches(5.05), Inches(4.67), Inches(0.3),
      font_name="Montserrat", font_size=9, color=SILVER,
      align=PP_ALIGN.CENTER)
txbox(s, "[ e.g. Content Creators, Photographers ]",
      Inches(2.0), Inches(5.35), Inches(4.67), Inches(0.42),
      font_name="Arial Black", font_size=14, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)

txbox(s, "OFFERING",
      Inches(6.67), Inches(5.05), Inches(4.67), Inches(0.3),
      font_name="Montserrat", font_size=9, color=SILVER,
      align=PP_ALIGN.CENTER)
txbox(s, "[ e.g. Brand Partnerships, Paid Gigs ]",
      Inches(6.67), Inches(5.35), Inches(4.67), Inches(0.42),
      font_name="Arial Black", font_size=14, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)

txbox(s, "@[HANDLE]  ·  [ businesswebsite.com ]",
      Inches(0.5), Inches(6.55), Inches(12.33), Inches(0.45),
      font_name="Montserrat", font_size=13, color=BLUE,
      align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════════════
out_path = "/Users/zoetaylor/Desktop/LINKDUP_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
